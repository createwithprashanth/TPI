"""
XYRA AI – EPC Sales Deck  (Black & White, v2)
Narrative arc: EPC Reality → AI Gap → XYRA Solution → Product → Pilot
Generates: XYRA_AI_EPC_Deck_v2.pptx
Run: python generate_xyra_epc_deck_v2.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

LOGO  = "/Users/prashanththipparthi/Desktop/xyra-ai/public/images/XYRA(PNG).png"
OUT   = Path(__file__).with_name("XYRA_AI_EPC_Deck_v2.pptx")
W     = Inches(13.333)
H     = Inches(7.5)
TOTAL = 11

# ── Palette ───────────────────────────────────────────────────────────────────
BLACK      = RGBColor(0,   0,   0)
NEAR_BLACK = RGBColor(18,  18,  18)
DARK_GRAY  = RGBColor(55,  55,  55)
MID_GRAY   = RGBColor(105, 105, 105)
LIGHT_GRAY = RGBColor(162, 162, 162)
RULE_GRAY  = RGBColor(210, 210, 210)
SOFT_BG    = RGBColor(245, 245, 245)
WHITE      = RGBColor(255, 255, 255)


# ── Primitives ─────────────────────────────────────────────────────────────────

def txt(slide, value, x, y, w, h,
        size=11, color=DARK_GRAY, bold=False,
        align=None, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.clear()
    for idx, line in enumerate(str(value).split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name   = "Arial"
        p.font.size   = Pt(size)
        p.font.bold   = bold
        p.font.italic = italic
        p.font.color.rgb = color
        if align is not None:
            p.alignment = align
    return box


def rect(slide, x, y, w, h, fill=WHITE, border=None, bw=Pt(0.8)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid();  s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border;  s.line.width = bw
    else:
        s.line.fill.background()
    return s


def hline(slide, x1, y, x2, color=RULE_GRAY, width=Pt(0.5)):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y, x2, y)
    c.line.color.rgb = color;  c.line.width = width


def vline(slide, x, y1, y2, color=RULE_GRAY, width=Pt(0.8)):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y1, x, y2)
    c.line.color.rgb = color;  c.line.width = width


def logo(slide, x, y, h=Inches(0.55)):
    try:
        slide.shapes.add_picture(LOGO, x, y, height=h)
    except Exception:
        pass


def footer(slide, n):
    hline(slide, Inches(0.55), Inches(7.08), Inches(12.78))
    txt(slide, "XYRA AI  ·  Private AI for EPC Engineering",
        Inches(0.55), Inches(7.13), Inches(6.0), Inches(0.22), 7.5, LIGHT_GRAY)
    txt(slide, f"{n} / {TOTAL}", Inches(12.2), Inches(7.13), Inches(0.9), Inches(0.22),
        7.5, LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def page_header(slide, eyebrow, title, subtitle=""):
    txt(slide, eyebrow.upper(), Inches(0.62), Inches(0.32), Inches(10), Inches(0.22),
        8.2, LIGHT_GRAY, True)
    txt(slide, title, Inches(0.58), Inches(0.6), Inches(10.5), Inches(0.72),
        25, BLACK, True)
    if subtitle:
        txt(slide, subtitle, Inches(0.62), Inches(1.36), Inches(10.8), Inches(0.36),
            11.5, MID_GRAY)
    hline(slide, Inches(0.58), Inches(1.82), Inches(12.78))


def card(slide, x, y, w, h, title, body, top_h=Inches(0.07)):
    rect(slide, x, y, w, h, SOFT_BG, RULE_GRAY)
    rect(slide, x, y, w, top_h, BLACK)
    txt(slide, title, x + Inches(0.2), y + Inches(0.19),
        w - Inches(0.4), Inches(0.3), 10.5, NEAR_BLACK, True)
    txt(slide, body,  x + Inches(0.2), y + Inches(0.57),
        w - Inches(0.4), h - Inches(0.68), 9, MID_GRAY)


def metric(slide, value, label, x, y):
    txt(slide, value, x, y, Inches(2.0), Inches(0.45), 24, BLACK, True)
    rect(slide, x, y + Inches(0.47), Inches(1.1), Inches(0.03), BLACK)
    txt(slide, label, x, y + Inches(0.56), Inches(2.4), Inches(0.28), 8, MID_GRAY)


def flow_steps(slide, steps, x0, y, gap):
    for i, (label, sub) in enumerate(steps):
        cx = x0 + i * gap
        circle = slide.shapes.add_shape(9, cx, y, Inches(0.65), Inches(0.65))
        if i == len(steps) - 1:
            circle.fill.solid(); circle.fill.fore_color.rgb = BLACK
        else:
            circle.fill.solid(); circle.fill.fore_color.rgb = SOFT_BG
        circle.line.color.rgb = BLACK;  circle.line.width = Pt(1.2)
        nc = WHITE if i == len(steps) - 1 else BLACK
        txt(slide, str(i + 1), cx + Inches(0.19), y + Inches(0.14),
            Inches(0.28), Inches(0.28), 12, nc, True, PP_ALIGN.CENTER)
        txt(slide, label, cx - Inches(0.15), y + Inches(0.78),
            Inches(0.96), Inches(0.28), 8.2, NEAR_BLACK, True, PP_ALIGN.CENTER)
        if sub:
            txt(slide, sub, cx - Inches(0.15), y + Inches(1.08),
                Inches(0.96), Inches(0.22), 7.2, MID_GRAY, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arr = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                cx + Inches(0.71), y + Inches(0.32),
                cx + gap - Inches(0.06), y + Inches(0.32))
            arr.line.color.rgb = RULE_GRAY;  arr.line.width = Pt(1.2)


# ── SLIDE 1: Cover ─────────────────────────────────────────────────────────────

def s1_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE

    rect(s, 0, 0, Inches(6.0), H, NEAR_BLACK)
    for i in range(7):
        vline(s, Inches(6.0 + i * 1.05), 0, H, RULE_GRAY, Pt(0.4))

    logo(s, Inches(0.55), Inches(0.65), h=Inches(0.70))

    txt(s, "XYRA AI", Inches(0.55), Inches(1.72), Inches(5.1), Inches(0.35), 11, LIGHT_GRAY, True)
    txt(s, "XYRA\nStudio", Inches(0.5), Inches(2.18), Inches(5.2), Inches(2.0), 52, WHITE, True)
    txt(s, "Private AI for EPC Drawing Deliverables",
        Inches(0.55), Inches(4.48), Inches(5.1), Inches(0.4), 13.5, LIGHT_GRAY)
    hline(s, Inches(0.55), Inches(5.08), Inches(5.6), LIGHT_GRAY, Pt(0.5))
    txt(s, "Instrumentation  ·  Piping MTO  ·  Drawing Review\nDeployed inside your network — zero cloud exposure.",
        Inches(0.55), Inches(5.24), Inches(5.1), Inches(0.75), 10, LIGHT_GRAY)

    rect(s, 0, H - Inches(0.55), Inches(6.0), Inches(0.55), BLACK)
    txt(s, "Confidential — prepared for EPC evaluation",
        Inches(0.25), H - Inches(0.42), Inches(5.5), Inches(0.26),
        8, LIGHT_GRAY, align=PP_ALIGN.CENTER)

    txt(s, "Built for EPC Teams", Inches(6.55), Inches(1.42), Inches(6.3), Inches(0.45), 21, BLACK, True)
    txt(s, "XYRA Studio turns P&IDs and engineering PDFs into structured,\nreview-ready deliverables — inside your client's secure network.",
        Inches(6.55), Inches(2.05), Inches(6.3), Inches(0.65), 11.5, DARK_GRAY)

    for i, (t, b) in enumerate([
        ("P&ID Intelligence",    "Auto-extract instrument tags, loops, and line numbers from vector and scanned drawings."),
        ("Piping MTO",           "Computer-vision component detection with size extraction and EPC-style Excel output."),
        ("Centralized Database", "Shared engineering records — instruments, loops, IO, and piping — across disciplines."),
        ("On-Premise Deployment","Docker. One port. Local LLM. Client drawings never leave the client network."),
    ]):
        card(s, Inches(6.55), Inches(3.05 + i * 0.97), Inches(6.3), Inches(0.82), t, b)


# ── SLIDE 2: The EPC Reality (Story Slide 1) ──────────────────────────────────

def s2_epc_reality(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE

    # Full-width dark hook banner
    rect(s, 0, 0, W, Inches(1.78), NEAR_BLACK)
    txt(s, "ON A TYPICAL EPC PROJECT, LESS THAN HALF OF ENGINEERING TIME IS SPENT ON ACTUAL ENGINEERING.",
        Inches(0.62), Inches(0.26), Inches(11.2), Inches(0.62), 18.5, WHITE, True)
    txt(s, "The rest — manual take-off, data re-entry across systems, reformatting, and re-checking — is clerical overhead that costs millions.",
        Inches(0.62), Inches(0.98), Inches(11.5), Inches(0.4), 11, LIGHT_GRAY)

    # 4 story-beat cards in 2×2 grid
    beats = [
        ("Clerical Work,\nNot Engineering",
         "An engineer spends 2–4 hours extracting data from a single P&ID.\n"
         "On a 500-drawing package, that is weeks of specialist time consumed before a single engineering decision is made."),
        ("The Same Data,\nEntered 4–6 Times",
         "The instrument tag FIC-1234 is entered into the Instrument Index, the PDMS, procurement software, "
         "and the control system — separately, manually, by different people, often in the same week."),
        ("Discipline to Discipline\nThrough Disconnected Files",
         "The piping team needs the updated line list. An email goes out with last week's spreadsheet. "
         "The instrumentation team made 12 corrections since then — none of which are visible to the recipient."),
        ("Checking What\nWas Already Checked",
         "Review meetings re-verify data that was verified in the previous phase. "
         "Every drawing revision restarts the cycle. The same engineering decision gets made two, three, four times."),
    ]
    for i, (title, body) in enumerate(beats):
        row, col = i // 2, i % 2
        card(s,
             Inches(0.62 + col * 6.22),
             Inches(1.98 + row * 2.08),
             Inches(6.02), Inches(1.88), title, body)

    # Punchline band
    rect(s, Inches(0.62), Inches(6.28), Inches(12.1), Inches(0.56), SOFT_BG, RULE_GRAY)
    txt(s,
        "\"Your most expensive resource — the engineer — is being used as a data-entry service. That changes with XYRA AI.\"",
        Inches(0.88), Inches(6.38), Inches(11.5), Inches(0.36),
        11.5, DARK_GRAY, italic=True, align=PP_ALIGN.CENTER)

    footer(s, 2)


# ── SLIDE 3: The AI Gap  (Story Slide 2) ─────────────────────────────────────

def s3_ai_gap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE

    # Dark left panel
    rect(s, 0, 0, Inches(6.22), H, NEAR_BLACK)

    # Left: Generic AI failure
    txt(s, "GENERIC AI", Inches(0.62), Inches(0.52), Inches(5.2), Inches(0.26), 8.5, LIGHT_GRAY, True)
    txt(s, "Why it fails\nEPC teams.", Inches(0.62), Inches(0.84), Inches(5.2), Inches(0.82), 23, WHITE, True)

    failures = [
        ("Not ISA-5.1 aware",
         "General LLMs don't understand instrument type decode, loop logic, or EPC tagging conventions."),
        ("Can't reliably read P&IDs",
         "ChatGPT and similar tools cannot extract structured engineering data from drawing PDFs with accuracy."),
        ("Client drawings can't leave the network",
         "Uploading proprietary P&IDs to a public AI service violates confidentiality on most EPC contracts."),
        ("Outputs can't be trusted",
         "Hallucinated instrument tags or wrong IO types are dangerous. There is no audit trail, no confidence score."),
        ("No project context",
         "Generic AI has no knowledge of your legend, naming conventions, line classes, or drawing standards."),
    ]
    DARK_CARD  = RGBColor(32, 32, 32)
    DARK_BORD  = RGBColor(58, 58, 58)
    for i, (title, body) in enumerate(failures):
        fy = Inches(2.05 + i * 0.88)
        rect(s, Inches(0.62), fy, Inches(5.3), Inches(0.74), DARK_CARD, DARK_BORD)
        txt(s, title, Inches(0.82), fy + Inches(0.1),  Inches(4.8), Inches(0.26), 10,  WHITE,      True)
        txt(s, body,  Inches(0.82), fy + Inches(0.38), Inches(4.8), Inches(0.28), 8.5, LIGHT_GRAY)

    # Right: XYRA solution
    txt(s, "XYRA AI", Inches(6.68), Inches(0.52), Inches(6.0), Inches(0.26), 8.5, MID_GRAY, True)
    txt(s, "Built for EPC\nfrom the ground up.", Inches(6.68), Inches(0.84), Inches(6.0), Inches(0.82), 23, BLACK, True)

    solutions = [
        ("ISA-5.1 built in",
         "XYRA models are trained with first-letter decode, qualifier logic, and EPC tagging conventions baked in."),
        ("Reads vector and scanned P&IDs",
         "Processes both vector PDFs and scanned drawings. OCR + LLM extraction tuned for engineering drawings."),
        ("100% on your network",
         "All models run locally via Ollama. Client drawings never leave the client-controlled server. No API calls out."),
        ("Engineer-reviewed outputs",
         "Every extraction is flagged with confidence scores and QA notes. Engineer approves before any delivery."),
        ("Project context aware",
         "Loads your project legend, naming rules, and scope to guide extraction — not generic defaults."),
    ]
    for i, (title, body) in enumerate(solutions):
        sy = Inches(2.05 + i * 0.88)
        rect(s, Inches(6.68), sy, Inches(6.1),    Inches(0.74), SOFT_BG, RULE_GRAY)
        rect(s, Inches(6.68), sy, Inches(0.06),   Inches(0.74), BLACK)
        txt(s, title, Inches(6.88), sy + Inches(0.1),  Inches(5.7), Inches(0.26), 10,  BLACK,    True)
        txt(s, body,  Inches(6.88), sy + Inches(0.38), Inches(5.7), Inches(0.28), 8.5, MID_GRAY)

    # CTA band (no standard footer — the band closes the slide)
    rect(s, 0, Inches(6.5), W, Inches(0.68), BLACK)
    txt(s, "One extraction. Structured outputs to every discipline. Inside your network.",
        Inches(0.62), Inches(6.66), Inches(12.1), Inches(0.38),
        15, WHITE, True, PP_ALIGN.CENTER)
    txt(s, f"3 / {TOTAL}", Inches(12.2), Inches(7.13), Inches(0.9), Inches(0.22),
        7.5, LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ── SLIDE 4: Platform Overview ────────────────────────────────────────────────

def s4_platform(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    rect(s, 0, 0, W, Inches(0.1), BLACK)
    logo(s, Inches(12.1), Inches(0.22), h=Inches(0.44))
    page_header(s, "XYRA Studio", "One platform. Four integrated tools.",
                "All tools share a single workspace. Upload once — use across every tool without re-uploading.")

    tools = [
        ("01", "InstruMap",    "Instrumentation Intelligence",
         "Extract tags · classify with ISA-5.1 logic · map to pipe lines\nOutputs: Instrument Index · IO List · Verification Log · Line List"),
        ("02", "Piping MTO",   "Component Take-Off",
         "User-built library + computer-vision detection across all drawings\nOutputs: MTO · Detection Register · QA Checks · Excel package"),
        ("03", "PrecisionPDF", "Drawing Review Workspace",
         "PDF viewer · annotation tools · text search · minimap · thumbnails\nSave and download annotated engineering drawings"),
        ("04", "System Health","Infrastructure Monitor",
         "Live status: API · Redis · Worker · Ollama · XYRA models\nClient-facing confidence that local compute is operational"),
    ]
    for i, (num, name, sub, body) in enumerate(tools):
        x = Inches(0.62 + i * 3.18)
        rect(s, x, Inches(2.0), Inches(2.98), Inches(4.35), SOFT_BG, RULE_GRAY)
        txt(s, num,  x + Inches(0.2), Inches(2.18), Inches(0.55), Inches(0.55), 26, RULE_GRAY, True)
        txt(s, name, x + Inches(0.2), Inches(2.8),  Inches(2.6),  Inches(0.38), 14.5, BLACK,   True)
        txt(s, sub,  x + Inches(0.2), Inches(3.26), Inches(2.6),  Inches(0.28), 9,  MID_GRAY, italic=True)
        hline(s, x + Inches(0.2), Inches(3.64), x + Inches(2.72))
        txt(s, body, x + Inches(0.2), Inches(3.8),  Inches(2.6),  Inches(2.0),  9,  DARK_GRAY)

    txt(s, "Shared workspace  ·  Background job processing  ·  Local LLM intelligence  ·  No re-upload required",
        Inches(1.5), Inches(6.6), Inches(10.3), Inches(0.24), 9, LIGHT_GRAY, align=PP_ALIGN.CENTER)
    footer(s, 4)


# ── SLIDE 5: InstruMap ────────────────────────────────────────────────────────

def s5_instrumap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    rect(s, 0, 0, W, Inches(0.1), BLACK)
    logo(s, Inches(12.1), Inches(0.22), h=Inches(0.44))
    page_header(s, "InstruMap", "Intelligent instrumentation extraction from P&IDs.",
                "Processes vector and scanned P&ID PDFs. Extracts, classifies, maps, and delivers EPC-ready outputs.")

    for i, (t, b) in enumerate([
        ("Tag Extraction",       "Detects instrument tags from vector text and OCR. Filters noise, equipment refs, and title-block fragments."),
        ("ISA-5.1 Classification","Applies first-letter and subsequent-letter decode. Resolves HH/LL qualifiers. Flags uncertain types for review."),
        ("Line Mapping",         "Maps each instrument to its connected pipe line using geometry and proximity. Infers service context."),
        ("Service Inference",    "Uses upstream/downstream context, EPC conventions, and project legend to infer instrument service description."),
    ]):
        card(s, Inches(0.62), Inches(2.05 + i * 1.28), Inches(4.55), Inches(1.12), t, b, Inches(0.06))

    txt(s, "DELIVERABLES GENERATED", Inches(5.55), Inches(2.0), Inches(7.5), Inches(0.26), 7.5, MID_GRAY, True)

    for i, (t, b) in enumerate([
        ("Instrument Index",  "All detected tags, loops, types, service, drawing reference, and QA flags."),
        ("IO List",           "AI / AO / DI / DO point list arranged for controls and automation review."),
        ("Verification Log",  "Full extraction trail: raw text, confidence, suppression reason, and review notes."),
        ("Line List",         "Pipe line numbers extracted from the drawing with instrument associations."),
    ]):
        bx, by, bw, bh = Inches(5.55), Inches(2.4 + i * 1.18), Inches(7.25), Inches(0.98)
        rect(s, bx, by, bw, bh, WHITE, RULE_GRAY)
        rect(s, bx, by, Inches(0.42), bh, NEAR_BLACK)
        txt(s, str(i + 1), bx + Inches(0.09), by + Inches(0.32), Inches(0.24), Inches(0.24), 10, WHITE, True, PP_ALIGN.CENTER)
        txt(s, t, bx + Inches(0.58), by + Inches(0.14), Inches(6.5), Inches(0.3), 11, BLACK, True)
        txt(s, b, bx + Inches(0.58), by + Inches(0.5),  Inches(6.5), Inches(0.36), 9, MID_GRAY)

    txt(s, "Delivered as a ZIP package — review-ready in Excel on open.",
        Inches(5.55), Inches(7.0), Inches(7.2), Inches(0.22), 8.5, LIGHT_GRAY, italic=True)
    footer(s, 5)


# ── SLIDE 6: Piping MTO ───────────────────────────────────────────────────────

def s6_mto(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    rect(s, 0, 0, W, Inches(0.1), BLACK)
    logo(s, Inches(12.1), Inches(0.22), h=Inches(0.44))
    page_header(s, "Piping MTO", "Computer-vision component detection across P&ID packages.",
                "User builds a component library once. XYRA detects and counts across every drawing and page.")

    flow_steps(s, [
        ("Upload\nP&IDs", ""),
        ("Select\nComponent", ""),
        ("Save to\nLibrary", ""),
        ("Run\nDetection", ""),
        ("Review\nResults", ""),
        ("Export\nPackage", ""),
    ], Inches(0.88), Inches(2.08), Inches(2.05))

    for i, (t, b) in enumerate([
        ("Piping MTO",        "Component counts per type with drawing and page reference."),
        ("Detection Register","Every match: page, location, size, confidence, and review status."),
        ("QA Checks",         "Automated checks on count thresholds and unreviewed detection items."),
        ("AI Review",         "Local XYRA MTO reviewer provides QA notes where the model is available."),
    ]):
        card(s, Inches(0.62 + i * 3.18), Inches(3.7), Inches(2.98), Inches(1.85), t, b)

    txt(s, "Size extraction  ·  Rotation support  ·  Tolerant match mode  ·  ORB fallback  ·  Multi-page PDF",
        Inches(1.5), Inches(6.65), Inches(10.3), Inches(0.24), 8.8, LIGHT_GRAY, align=PP_ALIGN.CENTER)
    footer(s, 6)


# ── SLIDE 7: Centralized Engineering Database ─────────────────────────────────

def s7_database(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    rect(s, 0, 0, W, Inches(0.1), BLACK)
    logo(s, Inches(12.1), Inches(0.22), h=Inches(0.44))
    page_header(s, "Centralized Engineering Database",
                "Discipline to discipline — connected, not copied.",
                "Structured PostgreSQL database. One record per project. Every discipline reads from the same source of truth.")

    # Problem → Solution contrast strip
    rect(s, Inches(0.62), Inches(1.95), Inches(5.9), Inches(0.82), SOFT_BG, RULE_GRAY)
    txt(s, "WITHOUT XYRA DB", Inches(0.82), Inches(2.02), Inches(5.5), Inches(0.22), 7.5, MID_GRAY, True)
    txt(s, "Instrument list in 6 different Excel files · piping quantities in emails · no version control · each discipline working from a different snapshot",
        Inches(0.82), Inches(2.26), Inches(5.5), Inches(0.42), 9, DARK_GRAY)

    rect(s, Inches(6.82), Inches(1.95), Inches(6.0), Inches(0.82), NEAR_BLACK)
    txt(s, "WITH XYRA DB", Inches(7.02), Inches(2.02), Inches(5.6), Inches(0.22), 7.5, LIGHT_GRAY, True)
    txt(s, "One database per project. Every discipline queries the same record. Revisions tracked. Outputs generated on demand.",
        Inches(7.02), Inches(2.26), Inches(5.6), Inches(0.42), 9, WHITE)

    hline(s, Inches(0.62), Inches(2.85), Inches(12.78))

    for i, (t, b) in enumerate([
        ("Instrument Records",
         "Every tag, loop, type, service, IO point, and drawing reference stored per project with full revision history."),
        ("Piping Components",
         "MTO library and detection results indexed by project, drawing, and revision for cross-run comparison."),
        ("Line Associations",
         "Instrument-to-line mapping linked to pipe line records for full cross-discipline traceability."),
        ("Revision History",
         "Each extraction run is versioned. Compare results across drawing revisions without separate spreadsheets."),
        ("Multi-User Access",
         "Multiple engineers query and review the same project database simultaneously — no file-locking, no lost edits."),
        ("AI Retrieval Ready",
         "Structured records as the foundation for engineer Q&A: 'list all transmitters on 6-inch steam lines'."),
    ]):
        r, c = i // 3, i % 3
        card(s, Inches(0.62 + c * 4.22), Inches(3.1 + r * 1.92), Inches(4.0), Inches(1.72), t, b)

    rect(s, Inches(0.62), Inches(6.68), Inches(5.5), Inches(0.38), NEAR_BLACK)
    txt(s, "PostgreSQL  ·  Replaces spreadsheet-per-project workflows",
        Inches(0.82), Inches(6.76), Inches(5.1), Inches(0.22), 8.8, WHITE, True)
    footer(s, 7)


# ── SLIDE 8: Deployment & Security ───────────────────────────────────────────

def s8_deploy(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    rect(s, 0, 0, W, Inches(0.1), BLACK)
    logo(s, Inches(12.1), Inches(0.22), h=Inches(0.44))
    page_header(s, "Deployment & Security", "Built for client-controlled environments.",
                "Single-server Docker deployment. Minimal network exposure. No cloud calls. No data egress.")

    arch = [
        ("Engineer Browser",     "Chrome or Edge on any machine in the client network"),
        ("nginx + React SPA",    "Only port 80 is exposed — everything else is hidden behind Docker networks"),
        ("FastAPI Backend",      "API layer — internal Docker network only, not host-exposed"),
        ("RQ Worker",            "Background job processing — internal Docker network only"),
        ("Ollama + XYRA Models", "Local LLM inference — no external API, no model calls outside the server"),
    ]
    for i, (comp, note) in enumerate(arch):
        by = Inches(2.1 + i * 0.95)
        is_top = i == 0
        rect(s, Inches(0.62), by, Inches(4.1), Inches(0.76),
             NEAR_BLACK if is_top else SOFT_BG, RULE_GRAY)
        tc = WHITE if is_top else BLACK
        nc = LIGHT_GRAY if is_top else MID_GRAY
        txt(s, comp, Inches(0.84), by + Inches(0.14), Inches(3.7), Inches(0.3), 11, tc, True)
        txt(s, note, Inches(0.84), by + Inches(0.46), Inches(3.7), Inches(0.22), 8.5, nc)
        if i < len(arch) - 1:
            vline(s, Inches(2.72), by + Inches(0.76), by + Inches(0.95), RULE_GRAY, Pt(1))

    for i, (t, b) in enumerate([
        ("Single exposed port",   "Port 80 only. Backend, Redis, Worker, and Ollama stay hidden in Docker internal networks."),
        ("Data stays local",      "P&IDs, extracted tags, and all outputs remain on the client-controlled server at all times."),
        ("Local model inference", "XYRA models run on Ollama. No API key. No external model calls. Ever."),
        ("Security headers",      "nginx: X-Frame-Options, CSP, rate limiting, dotfile blocking, server tokens off."),
        ("Monthly licensing",     "Subscription model with controlled deployment, update, and support cycles."),
    ]):
        card(s, Inches(5.35), Inches(2.05 + i * 1.04), Inches(7.6), Inches(0.9), t, b, Inches(0.05))

    footer(s, 8)


# ── SLIDE 9: About XYRA AI ────────────────────────────────────────────────────

def s9_company(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    rect(s, 0, 0, W, Inches(0.1), BLACK)
    logo(s, Inches(12.1), Inches(0.22), h=Inches(0.44))
    page_header(s, "About XYRA AI", "Engineering AI built for process industry.",
                "We build private, on-premise AI tools that integrate directly into EPC and plant engineering workflows.")

    for i, (title, body, tags) in enumerate([
        ("Domain-First AI",
         "XYRA models are trained with ISA-5.1 and EPC drawing conventions built in — not adapted from general-purpose chatbots.",
         "Instrumentation  ·  P&ID logic  ·  ISA-5.1"),
        ("Private by Design",
         "Client drawings, tags, and engineering data never leave the client network. All models run locally on Ollama.",
         "No cloud  ·  No API keys  ·  No data egress"),
        ("Engineer in Control",
         "XYRA accelerates first-pass extraction and surfaces uncertainty. Every output is designed for engineer review.",
         "QA flags  ·  Confidence scores  ·  Audit trail"),
        ("EPC-Ready Deployment",
         "Docker on Windows Server or Linux. One exposed port. Monthly licensing. Minimal IT footprint.",
         "Docker  ·  Port 80  ·  Windows / Linux"),
    ]):
        x = Inches(0.62 + i * 3.18)
        card(s, x, Inches(2.05), Inches(2.98), Inches(3.0), title, body)
        txt(s, tags, x + Inches(0.2), Inches(4.74), Inches(2.6), Inches(0.3), 8, LIGHT_GRAY, italic=True)

    hline(s, Inches(0.62), Inches(5.36), Inches(12.78))
    txt(s, "XYRA AI was founded by practitioners who spent years on EPC engineering deliverables.\nWe built the tool we needed — then made it available to the industry.",
        Inches(1.5), Inches(5.52), Inches(10.3), Inches(0.55), 13, DARK_GRAY, align=PP_ALIGN.CENTER)

    metric(s, "4",    "intelligent tools",      Inches(1.5),  Inches(6.28))
    metric(s, "6+",   "EPC deliverable types",  Inches(4.2),  Inches(6.28))
    metric(s, "100%", "on-premise deployment",  Inches(7.1),  Inches(6.28))
    metric(s, "0",    "cloud dependencies",     Inches(10.0), Inches(6.28))
    footer(s, 9)


# ── SLIDE 10: PrecisionPDF ────────────────────────────────────────────────────

def s10_pdf(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    rect(s, 0, 0, W, Inches(0.1), BLACK)
    logo(s, Inches(12.1), Inches(0.22), h=Inches(0.44))
    page_header(s, "PrecisionPDF", "Engineering drawing review and markup inside XYRA Studio.",
                "Full PDF workspace designed around the P&ID review workflow. No external tools needed.")

    for i, (t, b) in enumerate([
        ("High-Fidelity Rendering", "PDF.js rendering with zoom, pan, and page navigation. Works with large-format engineering drawings."),
        ("Text Search",             "Searchable text layer — find instrument tags, line numbers, or equipment notes directly on the drawing."),
        ("Annotation Tools",        "Freehand drawing, shape tools, and annotation sidebar for review markup and issue callouts."),
        ("Minimap & Thumbnails",    "Minimap overlay and thumbnail sidebar for navigation on large multi-page P&ID sets."),
        ("Save / Download",         "Export the annotated PDF directly from the browser. Review evidence stays with the drawing."),
        ("Shared Workspace",        "Same drawings available from InstruMap and Piping MTO — no re-upload between tools."),
    ]):
        r, c = i // 3, i % 3
        card(s, Inches(0.62 + c * 4.22), Inches(2.05 + r * 2.18), Inches(4.0), Inches(1.98), t, b)

    footer(s, 10)


# ── SLIDE 11: Pilot Proposal ──────────────────────────────────────────────────

def s11_pilot(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    rect(s, 0, 0, W, Inches(0.1), BLACK)
    logo(s, Inches(12.1), Inches(0.22), h=Inches(0.44))
    page_header(s, "Pilot Proposal", "Prove value on a real drawing package — in two weeks.",
                "A structured pilot gives your team a clear decision: time saved, output quality, and deployment fit.")

    phases = [
        ("Week 1", "Configure & Load", [
            "Deploy on client server or test VM",
            "Load project legend and title block context",
            "Configure instrument model and MTO library",
            "Run first extraction pass and calibrate output",
        ]),
        ("Week 2", "Validate & Compare", [
            "Process 20–50 representative P&IDs",
            "Compare XYRA output against manual workflow",
            "Review Instrument Index, IO List, and MTO package",
            "Identify gaps and tune detection settings",
        ]),
        ("Decision Pack", "Delivered at Close", [
            "Coverage summary with exception log",
            "Time-saving estimate vs. manual baseline",
            "Deployment checklist for production rollout",
            "Monthly license proposal",
        ]),
    ]
    for i, (phase, sub, points) in enumerate(phases):
        x     = Inches(0.62 + i * 4.22)
        is_last = i == 2
        rect(s, x, Inches(2.05), Inches(4.0), Inches(3.62),
             NEAR_BLACK if is_last else SOFT_BG, RULE_GRAY)
        tc  = WHITE      if is_last else BLACK
        sc  = LIGHT_GRAY if is_last else MID_GRAY
        btc = WHITE      if is_last else DARK_GRAY
        txt(s, phase.upper(), x + Inches(0.25), Inches(2.22), Inches(3.5), Inches(0.24), 8.5, sc, True)
        txt(s, sub,           x + Inches(0.25), Inches(2.5),  Inches(3.5), Inches(0.38), 14,  tc, True)
        hline(s, x + Inches(0.25), Inches(2.96), x + Inches(3.72),
              LIGHT_GRAY if is_last else RULE_GRAY)
        for j, pt in enumerate(points):
            txt(s, f"·  {pt}", x + Inches(0.25), Inches(3.12 + j * 0.52),
                Inches(3.5), Inches(0.42), 9.5, btc)

    rect(s, Inches(0.62), Inches(6.1), Inches(12.1), Inches(0.88), BLACK)
    txt(s, "Ready to run XYRA Studio on one of your P&ID packages?  →  prashanth.thipparthi@outlook.com",
        Inches(0.95), Inches(6.35), Inches(11.5), Inches(0.4),
        13, WHITE, True, PP_ALIGN.CENTER)

    footer(s, 11)


# ── Build ──────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    for fn in [s1_cover, s2_epc_reality, s3_ai_gap, s4_platform,
               s5_instrumap, s6_mto, s7_database, s8_deploy,
               s9_company, s10_pdf, s11_pilot]:
        fn(prs)

    prs.save(str(OUT))
    print(f"Saved → {OUT}")


if __name__ == "__main__":
    build()
