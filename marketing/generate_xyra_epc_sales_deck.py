from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("XYRA_Studio_EPC_Sales_Deck.pptx")

W = Inches(13.333)
H = Inches(7.5)

INK = RGBColor(10, 17, 31)
NAVY = RGBColor(16, 29, 48)
SLATE = RGBColor(70, 84, 105)
MUTED = RGBColor(105, 117, 138)
TEAL = RGBColor(0, 164, 156)
BLUE = RGBColor(44, 103, 218)
GREEN = RGBColor(33, 145, 93)
ORANGE = RGBColor(225, 130, 46)
LIGHT = RGBColor(246, 248, 251)
BORDER = RGBColor(217, 224, 234)
WHITE = RGBColor(255, 255, 255)


def text(slide, value, x, y, w, h, size=13, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = value
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align is not None:
        p.alignment = align
    return box


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    top = slide.shapes.add_shape(1, 0, 0, W, Inches(0.14))
    top.fill.solid()
    top.fill.fore_color.rgb = TEAL
    top.line.fill.background()

    for i in range(7):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(7.4 + i * 0.78),
            Inches(0.14),
            Inches(5.1 + i * 0.78),
            H,
        )
        line.line.color.rgb = RGBColor(239, 243, 248)
        line.line.width = Pt(0.6)


def footer(slide, n):
    text(slide, "XYRA Studio | Private AI for EPC drawing deliverables", Inches(0.55), Inches(7.08), Inches(5.0), Inches(0.2), 7.5, RGBColor(145, 155, 171))
    text(slide, f"{n}/5", Inches(12.25), Inches(7.08), Inches(0.55), Inches(0.2), 7.5, RGBColor(145, 155, 171), align=PP_ALIGN.RIGHT)


def header(slide, eyebrow, title, subtitle):
    text(slide, eyebrow.upper(), Inches(0.62), Inches(0.5), Inches(4.8), Inches(0.22), 8.5, TEAL, True)
    text(slide, title, Inches(0.58), Inches(0.8), Inches(8.7), Inches(0.72), 28, INK, True)
    text(slide, subtitle, Inches(0.62), Inches(1.48), Inches(8.2), Inches(0.42), 12.2, SLATE)


def pill(slide, label, x, y, w, color=TEAL):
    r = slide.shapes.add_shape(1, x, y, w, Inches(0.34))
    r.fill.solid()
    r.fill.fore_color.rgb = RGBColor(236, 253, 251) if color == TEAL else RGBColor(239, 246, 255)
    r.line.color.rgb = color
    r.line.width = Pt(0.8)
    text(slide, label, x + Inches(0.1), y + Inches(0.08), w - Inches(0.2), Inches(0.16), 8.5, color, True, PP_ALIGN.CENTER)


def card(slide, x, y, w, h, title, body, accent=TEAL):
    r = slide.shapes.add_shape(1, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = WHITE
    r.line.color.rgb = BORDER
    r.line.width = Pt(0.9)
    a = slide.shapes.add_shape(1, x, y, w, Inches(0.07))
    a.fill.solid()
    a.fill.fore_color.rgb = accent
    a.line.fill.background()
    text(slide, title, x + Inches(0.22), y + Inches(0.22), w - Inches(0.44), Inches(0.3), 12.0, INK, True)
    text(slide, body, x + Inches(0.22), y + Inches(0.68), w - Inches(0.44), h - Inches(0.86), 9.4, MUTED)


def metric(slide, value, label, x, y, color=TEAL):
    text(slide, value, x, y, Inches(1.25), Inches(0.42), 21, color, True)
    text(slide, label, x, y + Inches(0.43), Inches(1.85), Inches(0.25), 8.2, MUTED)


def window(slide, x, y, w, h, title):
    shell = slide.shapes.add_shape(1, x, y, w, h)
    shell.fill.solid()
    shell.fill.fore_color.rgb = RGBColor(12, 18, 31)
    shell.line.color.rgb = RGBColor(40, 54, 76)
    shell.line.width = Pt(1)
    bar = slide.shapes.add_shape(1, x, y, w, Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(22, 31, 49)
    bar.line.fill.background()
    for i, c in enumerate([RGBColor(255, 95, 87), RGBColor(255, 189, 46), RGBColor(40, 201, 64)]):
        d = slide.shapes.add_shape(1, x + Inches(0.12 + i * 0.14), y + Inches(0.17), Inches(0.07), Inches(0.07))
        d.fill.solid()
        d.fill.fore_color.rgb = c
        d.line.fill.background()
    text(slide, title, x + Inches(0.45), y + Inches(0.14), Inches(2.7), Inches(0.14), 8.0, WHITE, True)


def mini_table(slide, x, y, w, h):
    window(slide, x, y, w, h, "Instrument_Index.xlsx")
    headers = ["Instrument", "Loop", "Service", "IO", "Review"]
    rows = [
        ["FIT-1762", "F-1762", "Feed gas flow", "AI", ""],
        ["FCV-1762", "F-1762", "Flow control valve", "AO", ""],
        ["PIT-573", "P-573", "Vessel pressure", "AI", ""],
        ["ZSH-221", "Z-221", "Valve open limit", "DI", "QA"],
    ]
    col = [1.0, 0.82, 1.45, 0.4, 0.58]
    x0 = x + Inches(0.28)
    y0 = y + Inches(0.78)
    for i, htxt in enumerate(headers):
        cx = x0 + Inches(sum(col[:i]))
        r = slide.shapes.add_shape(1, cx, y0, Inches(col[i]), Inches(0.28))
        r.fill.solid()
        r.fill.fore_color.rgb = RGBColor(35, 50, 76)
        r.line.color.rgb = RGBColor(52, 68, 95)
        text(slide, htxt, cx + Inches(0.05), y0 + Inches(0.075), Inches(col[i] - 0.1), Inches(0.1), 6.5, WHITE, True)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cx = x0 + Inches(sum(col[:ci]))
            cy = y0 + Inches(0.28 + ri * 0.32)
            r = slide.shapes.add_shape(1, cx, cy, Inches(col[ci]), Inches(0.32))
            r.fill.solid()
            r.fill.fore_color.rgb = RGBColor(16, 24, 39) if ri % 2 else RGBColor(20, 30, 48)
            r.line.color.rgb = RGBColor(38, 52, 78)
            text(slide, val, cx + Inches(0.05), cy + Inches(0.09), Inches(col[ci] - 0.1), Inches(0.1), 6.3, RGBColor(216, 224, 238))


def flow(slide, labels, x, y, gap=1.58):
    for i, label in enumerate(labels):
        bx = x + Inches(i * gap)
        c = slide.shapes.add_shape(9, bx, y, Inches(0.72), Inches(0.72))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(237, 252, 250) if i < len(labels) - 1 else RGBColor(234, 245, 255)
        c.line.color.rgb = TEAL if i < len(labels) - 1 else BLUE
        text(slide, str(i + 1), bx + Inches(0.23), y + Inches(0.17), Inches(0.25), Inches(0.18), 13, TEAL if i < len(labels) - 1 else BLUE, True, PP_ALIGN.CENTER)
        text(slide, label, bx - Inches(0.25), y + Inches(0.84), Inches(1.22), Inches(0.28), 8.2, INK, True, PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, bx + Inches(0.78), y + Inches(0.36), bx + Inches(gap - 0.08), y + Inches(0.36))
            line.line.color.rgb = BORDER
            line.line.width = Pt(1.4)


def slide_1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    pill(s, "Built for EPC teams", Inches(0.62), Inches(0.62), Inches(1.55), TEAL)
    pill(s, "Private deployment", Inches(2.32), Inches(0.62), Inches(1.55), BLUE)
    text(s, "XYRA Studio", Inches(0.58), Inches(1.14), Inches(5.2), Inches(0.6), 31, INK, True)
    text(s, "Private AI for EPC drawing deliverables.", Inches(0.62), Inches(1.86), Inches(6.2), Inches(0.45), 17, SLATE)
    text(s, "Turn P&IDs and PDFs into review-ready engineering outputs.", Inches(0.62), Inches(2.7), Inches(6.2), Inches(1.2), 29, INK, True)
    text(s, "XYRA helps instrumentation and piping teams extract tags, line context, MTO counts, and QA evidence inside the client network.", Inches(0.66), Inches(4.07), Inches(5.75), Inches(0.7), 12.5, SLATE)
    metric(s, "P&ID", "instrument index + IO list", Inches(0.72), Inches(5.26), TEAL)
    metric(s, "MTO", "piping take-off support", Inches(2.58), Inches(5.26), BLUE)
    metric(s, "PDF", "review and markup workflow", Inches(4.42), Inches(5.26), GREEN)

    mini_table(s, Inches(7.12), Inches(1.1), Inches(5.3), Inches(3.15))
    card(s, Inches(7.12), Inches(4.62), Inches(1.62), Inches(1.1), "Extract", "Tags, loops, lines, and equipment context.", TEAL)
    card(s, Inches(8.94), Inches(4.62), Inches(1.62), Inches(1.1), "Review", "Confidence flags and QA evidence.", BLUE)
    card(s, Inches(10.76), Inches(4.62), Inches(1.62), Inches(1.1), "Deliver", "Excel packs your team can check.", GREEN)
    footer(s, 1)


def slide_2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Why now", "Manual drawing take-off is a capacity problem.", "EPC teams need faster first-pass deliverables without pushing client drawings outside the network.")
    card(s, Inches(0.78), Inches(2.2), Inches(2.75), Inches(1.55), "Manual extraction", "Engineers spend hours copying tags, loops, line numbers, and symbol counts from PDFs.", ORANGE)
    card(s, Inches(3.78), Inches(2.2), Inches(2.75), Inches(1.55), "Review drift", "Excel files, PDF markups, and review notes often separate from the drawing evidence.", BLUE)
    card(s, Inches(6.78), Inches(2.2), Inches(2.75), Inches(1.55), "Data restrictions", "Many clients cannot send drawings to public cloud AI services.", TEAL)
    card(s, Inches(9.78), Inches(2.2), Inches(2.75), Inches(1.55), "Small teams", "Specialist engineers are expensive to keep on repetitive take-off work.", GREEN)
    text(s, "XYRA is positioned as the first-pass engineering assistant: it accelerates extraction, makes uncertainty visible, and keeps the engineer in control.", Inches(1.0), Inches(4.55), Inches(11.3), Inches(0.72), 21, INK, True, PP_ALIGN.CENTER)
    metric(s, "1", "internal deployment surface", Inches(3.2), Inches(5.72), TEAL)
    metric(s, "3", "drawing intelligence tools", Inches(5.65), Inches(5.72), BLUE)
    metric(s, "6+", "review-ready outputs", Inches(8.1), Inches(5.72), GREEN)
    footer(s, 2)


def slide_3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "What clients receive", "Engineering deliverables, not a chatbot demo.", "Outputs are designed for review, filtering, traceability, and handover.")
    mini_table(s, Inches(0.75), Inches(2.0), Inches(5.15), Inches(3.15))
    outputs = [
        ("Instrument Index", "Instrument tags, loops, service text, drawing references, and review flags."),
        ("IO List", "AI, AO, DI, DO points arranged for controls and automation review."),
        ("Piping MTO", "Symbol counts and page-level evidence for piping material take-off support."),
        ("Verification Log", "Raw extraction trail, confidence, suppression reason, and QA notes."),
        ("PrecisionPDF Review", "PDF viewing, markups, measurements, and drawing review workspace."),
    ]
    for i, (title, body) in enumerate(outputs):
        card(s, Inches(6.35), Inches(1.72 + i * 0.82), Inches(5.75), Inches(0.62), title, body, [TEAL, BLUE, GREEN, ORANGE, TEAL][i])
    footer(s, 3)


def slide_4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Deployment story", "Runs where the drawings live.", "A practical architecture for client internal networks, monthly licensing, and controlled support.")
    flow(s, ["User browser", "nginx", "API", "Worker", "Local LLM"], Inches(1.02), Inches(2.32), gap=2.05)
    card(s, Inches(0.85), Inches(4.25), Inches(2.75), Inches(1.45), "One exposed port", "Browser traffic enters through port 80. Backend and model services stay hidden.", TEAL)
    card(s, Inches(3.92), Inches(4.25), Inches(2.75), Inches(1.45), "Client data stays local", "P&IDs and extracted data remain on the client-controlled server.", BLUE)
    card(s, Inches(6.99), Inches(4.25), Inches(2.75), Inches(1.45), "Project context", "Client legends, project scope, and naming rules can guide the extraction workflow.", GREEN)
    card(s, Inches(10.06), Inches(4.25), Inches(2.75), Inches(1.45), "Support-ready roadmap", "Health dashboard, logs, bundles, and redeploy checks reduce support burden.", ORANGE)
    footer(s, 4)


def slide_5(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Pilot offer", "Prove value on a real drawing package.", "A focused pilot gives the EPC a clear decision: time saved, output quality, and deployment fit.")
    card(s, Inches(0.8), Inches(2.0), Inches(3.55), Inches(2.35), "Week 1: configure", "Install inside client network\nLoad project legend and sample drawings\nRun first extraction and review calibration", TEAL)
    card(s, Inches(4.88), Inches(2.0), Inches(3.55), Inches(2.35), "Week 2: validate", "Process 20-50 representative P&IDs\nCompare against manual workflow\nReview Instrument Index, IO List, MTO, and PDF markups", BLUE)
    card(s, Inches(8.94), Inches(2.0), Inches(3.55), Inches(2.35), "Decision pack", "Coverage summary\nReview exceptions\nDeployment checklist\nMonthly license proposal", GREEN)
    r = s.shapes.add_shape(1, Inches(1.1), Inches(5.35), Inches(11.15), Inches(0.78))
    r.fill.solid()
    r.fill.fore_color.rgb = NAVY
    r.line.fill.background()
    text(s, "Next step: run XYRA on one EPC drawing package and review the generated deliverables together.", Inches(1.35), Inches(5.55), Inches(10.65), Inches(0.26), 14, WHITE, True, PP_ALIGN.CENTER)
    footer(s, 5)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    for make_slide in [slide_1, slide_2, slide_3, slide_4, slide_5]:
        make_slide(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
