"""
XYRA AI — EPC Sales Deck  ·  White Edition
White background · characters · workflows · real screenshots

Run: python3 generate_xyra_screenplay.py
"""
from __future__ import annotations
import struct
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paths ─────────────────────────────────────────────────────────────────────
SS   = Path("/Users/prashanththipparthi/Desktop/XYRA Studio/marketing/screenshots")
LOGO = "/Users/prashanththipparthi/Desktop/xyra-ai/public/images/XYRA(PNG).png"
OUT  = Path("/Users/prashanththipparthi/Desktop/XYRA Studio/marketing") / "XYRA_AI_Screenplay_Deck.pptx"

SYS_HEALTH = SS / "xyra system.jpeg"
CHECKPRINT = SS / "checkprint_from system_epc_level.jpeg"
INST_IDX   = SS / "Instrumenet_index_output.jpeg"
LINE_LST   = SS / "line_list_output.jpeg"
MTO_DETECT = SS / "mto_checkprint.jpeg"
MTO_OUT    = SS / "mto.jpeg"

W, H   = Inches(13.333), Inches(7.5)
TOTAL  = 13
LR     = 3.7215

# ═══════════════════════════════════════════════════════════
#  PALETTE  (white-background first)
# ═══════════════════════════════════════════════════════════
INK    = RGBColor(15,  23,  42)   # slate-900  — headlines
BODY   = RGBColor(51,  65,  85)   # slate-700  — body text
MID    = RGBColor(100, 116, 139)  # slate-500  — secondary
MUTED  = RGBColor(148, 163, 184)  # slate-400  — captions
RULE   = RGBColor(226, 232, 240)  # slate-200  — borders
SOFT   = RGBColor(248, 250, 252)  # slate-50   — card bg
BLUE   = RGBColor(37,  99,  235)  # blue-600   — primary
BLUE_L = RGBColor(219, 234, 254)  # blue-100   — light bg
BLUE_D = RGBColor(29,  78,  216)  # blue-700   — dark
TEAL   = RGBColor(13,  148, 136)  # teal-600
TEAL_L = RGBColor(204, 251, 241)  # teal-100
AMBER  = RGBColor(217, 119,   6)  # amber-600  — warning
AMB_L  = RGBColor(254, 243, 199)  # amber-100
RED    = RGBColor(220,  38,  38)  # red-600    — problem
RED_L  = RGBColor(254, 226, 226)  # red-100
GREEN  = RGBColor(5,   150, 105)  # emerald-600
GRN_L  = RGBColor(209, 250, 229)  # emerald-100
SLATE  = RGBColor(71,  85,  105)  # slate for figures
WHITE  = RGBColor(255, 255, 255)
BLACK  = RGBColor(0,   0,   0)

ML = Inches(0.72)   # left margin


# ═══════════════════════════════════════════════════════════
#  PRIMITIVES
# ═══════════════════════════════════════════════════════════

def wh(s):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE

def T(slide, txt, x, y, w, h,
      sz=11, c=BODY, bold=False, align=None, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    tf.clear()
    for i, line in enumerate(str(txt).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name   = "Arial"
        p.font.size   = Pt(sz)
        p.font.bold   = bold
        p.font.italic = italic
        p.font.color.rgb = c
        if align:
            p.alignment = align
    return tb

def R(slide, x, y, w, h, fill=WHITE, border=None, bw=Pt(0.8)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = bw
    else:       s.line.fill.background()
    return s

def Ov(slide, x, y, w, h, fill=BLUE, border=None, bw=Pt(0.8)):
    s = slide.shapes.add_shape(9, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = bw
    else:       s.line.fill.background()
    return s

def HL(slide, x1, y, x2, col=RULE, w=Pt(0.5)):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y, x2, y)
    c.line.color.rgb = col; c.line.width = w

def VL(slide, x, y1, y2, col=RULE, w=Pt(0.5)):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y1, x, y2)
    c.line.color.rgb = col; c.line.width = w

def DG(slide, x1, y1, x2, y2, col=RULE, w=Pt(0.5)):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = col; c.line.width = w

def logo_sm(slide, x=None, y=Inches(0.26), h=Inches(0.34)):
    lw = int(h * LR)
    lx = x if x is not None else ML
    try: slide.shapes.add_picture(LOGO, lx, y, height=h)
    except Exception: pass

def logo_c(slide, y, h=Inches(0.9)):
    lw = int(h * LR)
    try: slide.shapes.add_picture(LOGO, int((W - lw) / 2), y, height=h)
    except Exception: pass

def slide_chrome(slide, n, category=""):
    """Standard white-slide chrome: small logo, category, rule, page number."""
    logo_sm(slide)
    if category:
        T(slide, category.upper(), W - Inches(5.2), Inches(0.3), Inches(4.8), Inches(0.22),
          7, MUTED, align=PP_ALIGN.RIGHT)
    HL(slide, ML, Inches(7.1), W - Inches(0.55), RULE, Pt(0.5))
    T(slide, f"{n}  /  {TOTAL}", W - Inches(1.55), Inches(7.18), Inches(1.4), Inches(0.22),
      7.5, MUTED, align=PP_ALIGN.RIGHT)

def jpeg_dims(path):
    with open(str(path), "rb") as f:
        if f.read(2) != b"\xff\xd8":
            return None, None
        while True:
            b = f.read(1)
            while b != b"\xff":
                b = f.read(1)
            while b == b"\xff":
                b = f.read(1)
            m = b[0]
            if m in (0xC0, 0xC1, 0xC2, 0xC3):
                f.read(3)
                ih = struct.unpack(">H", f.read(2))[0]
                iw = struct.unpack(">H", f.read(2))[0]
                return iw, ih
            elif m == 0xD9:
                return None, None
            else:
                ln = struct.unpack(">H", f.read(2))[0]
                f.seek(ln - 2, 1)

def place_img(slide, path, x, y, avail_w, avail_h, shadow=True):
    iw, ih = jpeg_dims(path)
    if not iw:
        return
    ratio = iw / ih
    if int(avail_w / ratio) <= avail_h:
        w = avail_w
        h = int(avail_w / ratio)
    else:
        h = avail_h
        w = int(avail_h * ratio)
    ox = x + int((avail_w - w) / 2)
    oy = y + int((avail_h - h) / 2)
    if shadow:
        R(slide, ox + Inches(0.1), oy + Inches(0.1), w, h, RGBColor(180, 180, 180))
    slide.shapes.add_picture(str(path), ox, oy, w, h)

def full_w_img(slide, path, y_top):
    iw, ih = jpeg_dims(path)
    if not iw:
        return
    img_w = W
    img_h = int(img_w / (iw / ih))
    slide.shapes.add_picture(str(path), 0, y_top, img_w, img_h)

def ic_x(slide, x, y, s, col=AMBER):
    p = int(s * 0.15)
    DG(slide, x+p, y+p, x+s-p, y+s-p, col, Pt(2.2))
    DG(slide, x+s-p, y+p, x+p, y+s-p, col, Pt(2.2))

def ic_check(slide, x, y, s, col=GREEN):
    DG(slide, x, int(y+s*.52), int(x+s*.38), int(y+s), col, Pt(2.2))
    DG(slide, int(x+s*.38), int(y+s), int(x+s), y, col, Pt(2.2))

def pill(slide, x, y, label, fill=BLUE, tc=WHITE):
    pw = max(Inches(0.15) * len(label) + Inches(0.36), Inches(1.4))
    ph = Inches(0.36)
    R(slide, x, y, pw, ph, fill)
    T(slide, label, x + Inches(0.14), y + Inches(0.07),
      pw - Inches(0.28), ph - Inches(0.14), 8.5, tc, True, PP_ALIGN.CENTER)
    return pw


# ═══════════════════════════════════════════════════════════
#  CHARACTER BUILDERS
# ═══════════════════════════════════════════════════════════

def person(slide, cx, y, size, fill=BLUE, label=None):
    """Flat-style human figure. cx = horizontal center, y = top of head."""
    s    = size
    hs   = int(s * 0.28)   # head radius (oval half-size)
    bw   = int(s * 0.32)   # body half-width
    bh   = int(s * 0.38)   # body height
    lw   = int(s * 0.11)   # leg width
    lh   = int(s * 0.26)   # leg height
    neck = int(s * 0.03)   # neck gap

    Ov(slide, cx - hs, y, hs * 2, hs * 2, fill)                          # head
    R(slide,  cx - bw, y + hs*2 + neck, bw*2, bh, fill)                  # body
    R(slide,  cx - bw, y + hs*2 + neck + bh + int(s*0.02), lw*2, lh, fill) # left leg
    R(slide,  cx + bw - lw*2, y + hs*2 + neck + bh + int(s*0.02), lw*2, lh, fill) # right leg

    if label:
        T(slide, label,
          cx - Inches(0.75), y + hs*2 + neck + bh + lh + int(s*0.08),
          Inches(1.5), Inches(0.28), 8.5, MID, align=PP_ALIGN.CENTER)


def ai_robot(slide, cx, y, size, crossed=False):
    """Geometric AI robot figure."""
    s  = size
    hw = int(s * 0.55)   # head width
    hh = int(s * 0.4)    # head height
    bw = int(s * 0.68)   # body width
    bh = int(s * 0.44)   # body height
    ey = int(s * 0.1)    # eye size

    # Antenna
    VL(slide, cx, y, y + int(s * 0.14), SLATE, Pt(2.0))
    Ov(slide, cx - int(s*0.07), y - int(s*0.07), int(s*0.14), int(s*0.14), SLATE)

    head_y = y + int(s*0.14)
    R(slide, cx - hw//2, head_y, hw, hh, SOFT, RULE, Pt(1.5))           # head box
    Ov(slide, cx - int(s*0.2), head_y + int(hh*0.28), ey, ey, SLATE)    # eye L
    Ov(slide, cx + int(s*0.1), head_y + int(hh*0.28), ey, ey, SLATE)    # eye R
    T(slide, "AI", cx - hw//2, head_y + int(hh*0.58), hw, int(hh*0.35),
      8, SLATE, True, PP_ALIGN.CENTER)

    body_y = head_y + hh + int(s*0.04)
    R(slide, cx - bw//2, body_y, bw, bh, SOFT, RULE, Pt(1.5))           # body box

    if crossed:
        ic_x(slide, cx - hw//2 - Inches(0.04), head_y - Inches(0.06),
             hw + Inches(0.08), RED)


# ═══════════════════════════════════════════════════════════
#  WORKFLOW BUILDER
# ═══════════════════════════════════════════════════════════

def workflow(slide, steps, x0, y, sw, sh, gap, bad_idx=None, accent=BLUE):
    """
    Horizontal workflow diagram.
    steps : list of (title, subtitle)
    bad_idx : set of step indices to mark as problem steps
    """
    for i, (title, sub) in enumerate(steps):
        cx = x0 + i * (sw + gap)
        is_bad = bad_idx and i in bad_idx
        bg  = RED_L  if is_bad else SOFT
        bdr = RED    if is_bad else RULE

        R(slide, cx, y, sw, sh, bg, bdr, Pt(1.5))

        # Step number pill at top
        pill_s = Inches(0.3)
        pill_x = cx + int((sw - pill_s) / 2)
        pill_y = y - int(pill_s / 2)
        Ov(slide, pill_x, pill_y, pill_s, pill_s, RED if is_bad else accent)
        T(slide, str(i + 1), pill_x, pill_y + Inches(0.06),
          pill_s, pill_s * 7 // 10, 9, WHITE, True, PP_ALIGN.CENTER)

        # Title
        T(slide, title, cx + Inches(0.1), y + Inches(0.18),
          sw - Inches(0.2), Inches(0.3), 10, INK, True, PP_ALIGN.CENTER)
        # Subtitle
        if sub:
            T(slide, sub, cx + Inches(0.1), y + Inches(0.52),
              sw - Inches(0.2), Inches(0.42), 8.5, MID, align=PP_ALIGN.CENTER)

        # Problem badge
        if is_bad:
            ic_x(slide, cx + sw - Inches(0.3), y + Inches(0.08), Inches(0.24), RED)

        # Arrow connector
        if i < len(steps) - 1:
            ax1 = cx + sw + Inches(0.05)
            ax2 = cx + sw + gap - Inches(0.05)
            ay  = int(y + sh / 2)
            col = RED if is_bad else RULE
            HL(slide, ax1, ay, ax2, col, Pt(1.2))
            T(slide, "›", ax2 - Inches(0.16), int(ay - Inches(0.14)),
              Inches(0.18), Inches(0.26), 14, col, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
#  SLIDES
# ═══════════════════════════════════════════════════════════

# ── 1 · COVER ─────────────────────────────────────────────
def s01(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    wh(s)

    # Thin left accent bar
    R(s, 0, 0, Inches(0.06), H, BLUE)
    # Thin top bar
    R(s, 0, 0, W, Inches(0.06), BLUE)

    logo_c(s, Inches(1.35), h=Inches(0.9))

    T(s, "XYRA Studio",
      Inches(0.5), Inches(2.58), W - Inches(1), Inches(0.95),
      52, INK, True, PP_ALIGN.CENTER)

    HL(s, Inches(4.2), Inches(3.72), Inches(9.13), BLUE, Pt(2.0))

    T(s, "Private AI for EPC Engineering",
      Inches(0.5), Inches(3.88), W - Inches(1), Inches(0.46),
      16, BLUE, align=PP_ALIGN.CENTER)

    # Four capability chips centred
    chips = ["P&ID Intelligence", "Piping MTO", "Centralized Database", "On-Premise"]
    total_w = sum(max(Inches(0.15)*len(c)+Inches(0.36), Inches(1.4)) + Inches(0.3) for c in chips)
    cx = int((W - total_w) / 2)
    for chip in chips:
        cw = pill(s, cx, Inches(4.72), chip, BLUE, WHITE)
        cx += cw + Inches(0.3)

    T(s, "P&ID extraction  ·  Piping MTO  ·  PDF review  ·  Centralized project database",
      Inches(0.5), Inches(5.42), W - Inches(1), Inches(0.3),
      10, MUTED, align=PP_ALIGN.CENTER)

    # Bottom metadata bar
    R(s, 0, H - Inches(0.62), W, Inches(0.62), SOFT)
    HL(s, 0, H - Inches(0.62), W, RULE, Pt(0.6))
    T(s, "Prepared for  [ EPC Client ]  ·  2026",
      ML, H - Inches(0.44), Inches(7), Inches(0.26), 8.5, MUTED)
    T(s, "Confidential",
      W - Inches(2.5), H - Inches(0.44), Inches(2.3), Inches(0.26),
      8.5, MUTED, align=PP_ALIGN.RIGHT)


# ── 2 · MEET THE TEAM ─────────────────────────────────────
def s02(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    wh(s)
    slide_chrome(s, 2, "The People")

    T(s, "Meet the people XYRA was built for.",
      ML, Inches(0.75), Inches(10), Inches(0.68), 28, INK, True)
    T(s, "Instrumentation and piping engineers who spend most of their day on data — not decisions.",
      ML, Inches(1.5), Inches(11), Inches(0.36), 13, MID)

    HL(s, ML, Inches(2.0), W - ML, RULE, Pt(0.5))

    # Two character figures
    chars = [
        (Inches(1.62), BLUE,  "Instrumentation\nEngineer"),
        (Inches(3.42), TEAL,  "Piping\nEngineer"),
    ]
    for cx, col, lbl in chars:
        person(s, cx, Inches(2.52), Inches(1.55), col, lbl)

    # What they deal with — time breakdown bars
    bx = Inches(5.1)
    bw_max = Inches(6.8)
    T(s, "WHERE THE HOURS ACTUALLY GO", bx, Inches(2.28), Inches(7.5), Inches(0.24), 7.5, MUTED, True)

    tasks = [
        ("Manual P&ID take-off",              3.5, AMBER),
        ("Data re-entry across systems",       2.5, AMBER),
        ("Cross-discipline file transfers",    1.5, RED),
        ("Review & re-checking same data",     1.2, RED),
        ("Actual engineering decisions",       0.8, GREEN),
    ]
    max_val = max(t[1] for t in tasks)
    for i, (label, hrs, col) in enumerate(tasks):
        row_y = Inches(2.62 + i * 0.76)
        T(s, label, bx, row_y, Inches(3.2), Inches(0.28), 9.5, BODY)
        bar_w = int(bw_max * hrs / max_val)
        R(s, bx + Inches(3.32), row_y + Inches(0.04), bar_w, Inches(0.28),
          GRN_L if col == GREEN else (AMB_L if col == AMBER else RED_L), col, Pt(1.0))
        T(s, f"{hrs} hrs", bx + Inches(3.32) + bar_w + Inches(0.1), row_y,
          Inches(0.7), Inches(0.28), 9, MID)

    HL(s, bx, Inches(6.44), bx + bw_max + Inches(0.8), RULE, Pt(0.5))
    T(s, "Less than 1 hour of an 8-hour workday is spent on actual engineering decisions.",
      bx, Inches(6.55), Inches(7.8), Inches(0.32), 11, BLUE, True)

    slide_chrome(s, 2, "The People")


# ── 3 · THE BROKEN WORKFLOW ───────────────────────────────
def s03(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    wh(s)
    slide_chrome(s, 3, "The Problem")

    T(s, "This is how engineering data moves today.",
      ML, Inches(0.75), Inches(11), Inches(0.65), 28, INK, True)
    T(s, "From the moment a P&ID is issued to the moment a deliverable is produced.",
      ML, Inches(1.48), Inches(11), Inches(0.34), 13, MID)

    # Broken workflow — 6 steps
    bsteps = [
        ("P&ID\nIssued",          "Drawing arrives"),
        ("Manual\nExtraction",    "2–4 hrs per sheet"),
        ("Spreadsheet\nCreated",  "System A"),
        ("Email to\nPiping Team", "Last week's file"),
        ("Re-entered\ninto PDMS", "System B — manual"),
        ("Review\nMeeting",       "Checks what\nwas checked"),
    ]
    sw, sh, gap = Inches(1.72), Inches(0.9), Inches(0.29)
    total = len(bsteps) * sw + (len(bsteps)-1) * gap
    x0 = int((W - total) / 2)
    workflow(s, bsteps, x0, Inches(2.12), sw, sh, gap, bad_idx={1, 3, 4, 5}, accent=BLUE)

    # Problem callouts below the broken steps
    problems = [
        (1, "2–4 hrs of\nspecialist time"),
        (3, "Wrong version\nsent"),
        (4, "Re-keyed\nmanually"),
        (5, "Already checked\nlast phase"),
    ]
    for idx, note in problems:
        px = x0 + idx * (sw + gap) + int(sw / 2) - int(Inches(0.8))
        T(s, note, px, Inches(3.22), Inches(1.6), Inches(0.42), 7.5, RED, align=PP_ALIGN.CENTER)

    # The key insight
    R(s, ML, Inches(4.28), W - ML - Inches(0.55), Inches(0.78), AMB_L, AMBER, Pt(1.2))
    R(s, ML, Inches(4.28), Inches(0.05), Inches(0.78), AMBER)
    T(s, "The same instrument tag is entered 4–6 times across a single project.\n"
        "Discipline-to-discipline data transfer happens through disconnected files — no traceability, no version control.",
      ML + Inches(0.2), Inches(4.38), W - ML - Inches(0.88), Inches(0.58), 10.5, BODY)

    # Generic AI failure strip
    R(s, ML, Inches(5.25), W - ML - Inches(0.55), Inches(1.38), SOFT, RULE, Pt(1.0))
    ai_robot(s, Inches(1.45), Inches(5.38), Inches(0.92), crossed=True)
    T(s, "\"We'll use ChatGPT.\"", Inches(2.55), Inches(5.38), Inches(4.5), Inches(0.36), 15, INK, True, italic=True)
    fails = ["Doesn't know ISA-5.1  ✕",
             "Can't reliably read P&IDs  ✕",
             "Client drawings can't go to a public AI  ✕"]
    for i, f in enumerate(fails):
        T(s, f, Inches(2.55), Inches(5.82 + i * 0.26), Inches(9.0), Inches(0.26), 9.5, RED)

    slide_chrome(s, 3, "The Problem")


# ── 4 · THE COST ──────────────────────────────────────────
def s04(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    wh(s)
    slide_chrome(s, 4, "The Cost")

    # Left: big stat panel
    R(s, 0, 0, Inches(5.8), H, BLUE)
    T(s, "30–50%",
      Inches(0.3), Inches(1.6), Inches(5.2), Inches(2.4),
      80, WHITE, True, PP_ALIGN.CENTER)
    T(s, "of your engineering budget\ngoes to clerical work.",
      Inches(0.3), Inches(4.2), Inches(5.2), Inches(0.9),
      18, RGBColor(147, 197, 253), align=PP_ALIGN.CENTER)
    HL(s, Inches(0.6), Inches(5.28), Inches(5.2), RGBColor(29, 78, 216), Pt(0.8))
    T(s, "Not engineering.",
      Inches(0.3), Inches(5.44), Inches(5.2), Inches(0.52),
      22, WHITE, True, PP_ALIGN.CENTER)

    # Right: breakdown cards
    rx = Inches(6.3)
    rw = W - rx - Inches(0.55)

    T(s, "Where it goes", rx, Inches(0.82), rw, Inches(0.38), 18, INK, True)
    T(s, "Breakdown of wasted engineering capacity on a typical EPC project:",
      rx, Inches(1.26), rw, Inches(0.32), 11, MID)

    breakdown = [
        ("Manual take-off",           "2–4 hrs per P&ID drawing, repeated across hundreds of drawings",    AMBER,  AMB_L),
        ("Data re-entry",             "Same instrument tag entered 4–6× across systems by different teams", RED,    RED_L),
        ("Discipline handoffs",       "Emailed spreadsheets, wrong versions, no audit trail between teams",  AMBER,  AMB_L),
        ("Review & re-checking",      "Review phases repeat checks already done — no shared engineering record", RED, RED_L),
    ]
    for i, (title, body, col, bg) in enumerate(breakdown):
        cy = Inches(1.78 + i * 1.22)
        R(s, rx, cy, rw, Inches(1.08), bg, col, Pt(1.2))
        R(s, rx, cy, Inches(0.05), Inches(1.08), col)
        T(s, title, rx + Inches(0.2), cy + Inches(0.1), rw - Inches(0.28), Inches(0.28), 11, INK, True)
        T(s, body,  rx + Inches(0.2), cy + Inches(0.44), rw - Inches(0.28), Inches(0.52), 9, BODY)

    slide_chrome(s, 4, "The Cost")


# ── 5 · XYRA ENTERS ───────────────────────────────────────
def s05(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    wh(s)
    slide_chrome(s, 5, "The Solution")

    T(s, "There is a better way.",
      ML, Inches(0.78), Inches(11), Inches(0.72), 32, INK, True)

    HL(s, ML, Inches(1.6), ML + Inches(4.0), BLUE, Pt(1.5))

    logo_c(s, Inches(1.9), h=Inches(0.82))

    T(s, "XYRA Studio",
      Inches(0.5), Inches(2.95), W - Inches(1), Inches(0.6), 26, INK, True, PP_ALIGN.CENTER)
    T(s, "Private AI for EPC Drawing Deliverables",
      Inches(0.5), Inches(3.62), W - Inches(1), Inches(0.36), 14, BLUE, align=PP_ALIGN.CENTER)

    # 4 tool cards
    tools = [
        ("InstruMap",    "Extract instrument tags, classify with ISA-5.1, map to lines, generate Index / IO List / Verification Log.",  BLUE,  BLUE_L),
        ("Piping MTO",   "User-built component library + computer-vision detection across all drawings. EPC-style Excel output.",         TEAL,  TEAL_L),
        ("PrecisionPDF", "Full PDF viewer with annotation tools, text search, minimap, and thumbnail navigation built in.",              SLATE, SOFT),
        ("Centralized DB","PostgreSQL-backed engineering records — instruments, loops, piping, IO — shared across all disciplines.",     BLUE_D, BLUE_L),
    ]
    tw = Inches(2.82)
    tgap = Inches(0.26)
    tx = int((W - (4*tw + 3*tgap)) / 2)
    for i, (name, desc, col, bg) in enumerate(tools):
        cx = tx + i * (tw + tgap)
        R(s, cx, Inches(4.32), tw, Inches(2.38), bg, col, Pt(1.2))
        R(s, cx, Inches(4.32), tw, Inches(0.06), col)
        T(s, f"0{i+1}", cx + Inches(0.14), Inches(4.42), Inches(0.4), Inches(0.3), 11, col, True)
        T(s, name, cx + Inches(0.14), Inches(4.76), tw - Inches(0.28), Inches(0.32), 12, INK, True)
        T(s, desc, cx + Inches(0.14), Inches(5.14), tw - Inches(0.28), Inches(1.38), 9, BODY)

    slide_chrome(s, 5, "The Solution")


# ── 6 · THE XYRA WORKFLOW ─────────────────────────────────
def s06(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    wh(s)
    slide_chrome(s, 6, "How It Works")

    T(s, "This is how it works with XYRA.",
      ML, Inches(0.78), Inches(11), Inches(0.65), 28, INK, True)
    T(s, "One upload. Every tool. One structured output to all disciplines.",
      ML, Inches(1.5), Inches(11), Inches(0.34), 13, MID)

    HL(s, ML, Inches(1.98), W - ML, RULE, Pt(0.5))

    # Clean 4-step XYRA workflow
    xsteps = [
        ("Upload P&ID",     "Drag & drop\nPDF — any quality"),
        ("XYRA Extracts",   "ISA-5.1 decode\nLine mapping\nService inference"),
        ("Engineer Reviews","Confidence flags\nQA notes\nOne-click remove"),
        ("Download Outputs","Index · IO List\nMTO · Line List\nZIP package"),
    ]
    sw, sh, gap = Inches(2.55), Inches(1.35), Inches(0.52)
    total = 4*sw + 3*gap
    x0 = int((W - total) / 2)
    workflow(s, xsteps, x0, Inches(2.25), sw, sh, gap, bad_idx=None, accent=BLUE)

    # Confirmation icons below each step
    confirms = ["Drop PDFs\nno prep needed", "LLM + OCR\nruns locally", "Full audit trail\nQA flagged", "Excel + ZIP\nreview-ready"]
    for i, txt in enumerate(confirms):
        cx = x0 + i * (sw + gap)
        ic_check(s, cx + int((sw - Inches(0.3))/2), Inches(3.8), Inches(0.3), GREEN)
        T(s, txt, cx, Inches(4.2), sw, Inches(0.4), 8.5, MID, align=PP_ALIGN.CENTER)

    # Contrast with the old broken workflow
    R(s, ML, Inches(5.08), W - ML - Inches(0.55), Inches(1.72), SOFT, RULE, Pt(1.0))

    T(s, "BEFORE XYRA", ML + Inches(0.22), Inches(5.22), Inches(5.0), Inches(0.26), 7.5, RED, True)
    T(s, "Manual take-off  →  5 systems re-entered  →  email  →  wrong version  →  re-check",
      ML + Inches(0.22), Inches(5.52), Inches(5.3), Inches(0.56), 9.5, BODY, italic=True)

    VL(s, Inches(6.45), Inches(5.15), Inches(6.88), RULE, Pt(0.6))

    T(s, "WITH XYRA", Inches(6.62), Inches(5.22), Inches(5.0), Inches(0.26), 7.5, BLUE, True)
    T(s, "Upload  →  XYRA reads  →  engineer reviews  →  structured outputs to all disciplines",
      Inches(6.62), Inches(5.52), Inches(5.9), Inches(0.56), 9.5, BODY, italic=True)

    slide_chrome(s, 6, "How It Works")


# ── 7 · SYSTEM HEALTH ─────────────────────────────────────
def s07(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    wh(s)

    logo_sm(s)
    T(s, "XYRA STUDIO  ·  LIVE SYSTEM", W - Inches(5.2), Inches(0.3), Inches(4.8), Inches(0.22),
      7, MUTED, align=PP_ALIGN.RIGHT)

    T(s, "Compute fabric online.",
      ML, Inches(0.72), Inches(9), Inches(0.65), 30, INK, True)
    T(s, "8 / 8 services live  ·  4 custom AI engines  ·  18.7 GB models loaded  ·  100% health  ·  0 queue",
      ML, Inches(1.42), Inches(12), Inches(0.34), 12, MID)

    # Full-width screenshot anchored to bottom
    iw, ih = jpeg_dims(SYS_HEALTH)
    img_h = int(W / (iw / ih))
    img_y = H - img_h
    s.shapes.add_picture(str(SYS_HEALTH), 0, img_y, W, img_h)

    HL(s, ML, Inches(7.1), W - Inches(0.55), RULE, Pt(0.5))
    T(s, f"7  /  {TOTAL}", W - Inches(1.55), Inches(7.18), Inches(1.4), Inches(0.22),
      7.5, MUTED, align=PP_ALIGN.RIGHT)


# ── 8 · P&ID CHECKPRINT ───────────────────────────────────
def s08(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    wh(s)

    logo_sm(s)
    T(s, "P&ID INTELLIGENCE", W - Inches(5.2), Inches(0.3), Inches(4.8), Inches(0.22),
      7, MUTED, align=PP_ALIGN.RIGHT)

    T(s, "XYRA reads the drawing.",
      ML, Inches(0.72), Inches(9.5), Inches(0.65), 30, INK, True)
    T(s, "Every instrument tag identified, classified, and linked to its pipe line number.",
      ML, Inches(1.42), Inches(11), Inches(0.34), 12, MID)

    # P&ID screenshot — centred, letterboxed with a white background showing
    place_img(s, CHECKPRINT, Inches(0), Inches(1.88), W, H - Inches(1.88), shadow=False)

    HL(s, ML, Inches(7.1), W - Inches(0.55), RULE, Pt(0.5))
    T(s, f"8  /  {TOTAL}", W - Inches(1.55), Inches(7.18), Inches(1.4), Inches(0.22),
      7.5, MUTED, align=PP_ALIGN.RIGHT)


# ── 9 · INSTRUMENT INDEX ──────────────────────────────────
def s09(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    wh(s)

    logo_sm(s)
    T(s, "INSTRUMENTATION OUTPUT", W - Inches(5.2), Inches(0.3), Inches(4.8), Inches(0.22),
      7, MUTED, align=PP_ALIGN.RIGHT)

    T(s, "Every instrument.", ML, Inches(0.72), Inches(7), Inches(0.65), 30, INK, True)
    T(s, "Structured.", ML + Inches(5.4), Inches(0.72), Inches(5), Inches(0.65), 30, BLUE, True)
    T(s, "Instrument Index — tags · loops · ISA type · service description · IO type · P&ID reference · QA review flags",
      ML, Inches(1.42), Inches(12), Inches(0.34), 12, MID)

    full_w_img(s, INST_IDX, Inches(1.85))

    HL(s, ML, Inches(7.1), W - Inches(0.55), RULE, Pt(0.5))
    T(s, f"9  /  {TOTAL}", W - Inches(1.55), Inches(7.18), Inches(1.4), Inches(0.22),
      7.5, MUTED, align=PP_ALIGN.RIGHT)


# ── 10 · LINE LIST ────────────────────────────────────────
def s10(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    wh(s)

    logo_sm(s)
    T(s, "PIPING OUTPUT", W - Inches(5.2), Inches(0.3), Inches(4.8), Inches(0.22),
      7, MUTED, align=PP_ALIGN.RIGHT)

    T(s, "Every line number.", ML, Inches(0.72), Inches(8), Inches(0.65), 30, INK, True)
    T(s, "Mapped.", ML + Inches(6.18), Inches(0.72), Inches(5), Inches(0.65), 30, TEAL, True)
    T(s, "Line List — pipe numbers · sizes · fluid codes · insulation spec · area codes · coordinates · P&ID source",
      ML, Inches(1.42), Inches(12), Inches(0.34), 12, MID)

    full_w_img(s, LINE_LST, Inches(1.85))

    HL(s, ML, Inches(7.1), W - Inches(0.55), RULE, Pt(0.5))
    T(s, f"10  /  {TOTAL}", W - Inches(1.55), Inches(7.18), Inches(1.4), Inches(0.22),
      7.5, MUTED, align=PP_ALIGN.RIGHT)


# ── 11 · PIPING MTO ───────────────────────────────────────
def s11(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    wh(s)
    slide_chrome(s, 11, "Piping MTO")

    T(s, "Select once.", ML, Inches(0.78), Inches(9), Inches(0.65), 28, INK, True)
    T(s, "XYRA finds every instance — across every page, every rotation, with size extraction.",
      ML, Inches(1.5), Inches(11), Inches(0.34), 13, MID)

    # MTO workflow
    msteps = [
        ("Select\nComponent",  "Draw a box\naround any symbol"),
        ("Save to\nLibrary",   "Auto-trimmed\nthumbnail stored"),
        ("Run\nDetection",     "Computer vision\nacross all pages"),
        ("Review\nResults",    "Remove false\npositives"),
        ("Export\nPackage",    "MTO · Register\nQA · Excel ZIP"),
    ]
    sw, sh, gap = Inches(2.0), Inches(0.88), Inches(0.32)
    total = 5*sw + 4*gap
    x0 = int((W - total) / 2)
    workflow(s, msteps, x0, Inches(2.22), sw, sh, gap, bad_idx=None, accent=TEAL)

    # MTO detection screenshot
    place_img(s, MTO_DETECT, ML, Inches(3.42), Inches(7.5), Inches(3.3), shadow=True)

    # MTO output stats — right side
    rx = ML + Inches(7.8)
    rw = W - rx - Inches(0.55)

    T(s, "OUTPUT", rx, Inches(3.45), rw, Inches(0.24), 7.5, MUTED, True)

    place_img(s, MTO_OUT, rx, Inches(3.72), rw, Inches(1.5), shadow=True)

    stats_row = [("17", "components detected"), ("1", "drawing processed"), ("1", "component type")]
    for i, (val, lbl) in enumerate(stats_row):
        cy = Inches(5.45 + i * 0.52)
        T(s, val, rx, cy, Inches(0.6), Inches(0.4), 22, TEAL, True)
        T(s, lbl, rx + Inches(0.7), cy + Inches(0.05), rw - Inches(0.8), Inches(0.3), 9, MID)

    slide_chrome(s, 11, "Piping MTO")


# ── 12 · PILOT PLAN ───────────────────────────────────────
def s12(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    wh(s)
    slide_chrome(s, 12, "The Pilot")

    T(s, "Two weeks. One answer.",
      ML, Inches(0.78), Inches(11), Inches(0.65), 28, INK, True)
    T(s, "A structured pilot gives your team a clear decision: time saved, output quality, and deployment fit.",
      ML, Inches(1.5), Inches(11), Inches(0.34), 13, MID)

    HL(s, ML, Inches(1.98), W - ML, RULE, Pt(0.5))

    # 3 pilot phases as large workflow cards with character figures
    phases = [
        ("WEEK 1",        "Configure & Load",
         ["Deploy on client server or test VM",
          "Load project legend + title block context",
          "Configure instrument model + MTO library",
          "Run first extraction pass — calibrate output"],
         BLUE,  BLUE_L, Inches(1.8)),
        ("WEEK 2",        "Validate & Compare",
         ["Process 20–50 representative P&IDs",
          "Compare XYRA output vs manual workflow",
          "Review Instrument Index, IO List, MTO",
          "Identify gaps — tune detection settings"],
         TEAL,  TEAL_L, Inches(5.25)),
        ("DECISION PACK", "Delivered at Close",
         ["Coverage summary + exception log",
          "Time-saving estimate vs manual baseline",
          "Deployment checklist for production",
          "Monthly license proposal"],
         BLUE_D, BLUE_L, Inches(8.7)),
    ]
    pw = Inches(3.88)
    for phase, sub, points, col, bg, px in phases:
        R(s, px, Inches(2.22), pw, Inches(4.52), bg, col, Pt(1.5))
        R(s, px, Inches(2.22), pw, Inches(0.06), col)

        # Character figure for each phase
        person(s, int(px + pw/2), Inches(2.42), Inches(0.82), col)

        T(s, phase, px + Inches(0.2), Inches(3.55), pw - Inches(0.4), Inches(0.24), 7.5, col, True)
        T(s, sub,   px + Inches(0.2), Inches(3.82), pw - Inches(0.4), Inches(0.36), 14, INK, True)
        HL(s, px + Inches(0.2), Inches(4.28), px + pw - Inches(0.2), RULE, Pt(0.5))
        for i, pt in enumerate(points):
            ic_check(s, px + Inches(0.2), Inches(4.44 + i * 0.48), Inches(0.24), col)
            T(s, pt, px + Inches(0.58), Inches(4.42 + i * 0.48), pw - Inches(0.7), Inches(0.38), 9.5, BODY)

    slide_chrome(s, 12, "The Pilot")


# ── 13 · CLOSE ────────────────────────────────────────────
def s13(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    wh(s)

    # Thin accent bars
    R(s, 0, 0, Inches(0.06), H, BLUE)
    R(s, 0, 0, W, Inches(0.06), BLUE)

    logo_c(s, Inches(1.2), h=Inches(0.72))

    T(s, "Ready to see XYRA\non your drawings?",
      Inches(0.5), Inches(2.32), W - Inches(1), Inches(2.0),
      46, INK, True, PP_ALIGN.CENTER)

    HL(s, Inches(4.2), Inches(4.52), Inches(9.13), BLUE, Pt(1.5))

    T(s, "We deploy inside your network, run XYRA on your drawing package,\n"
        "and deliver structured outputs your team evaluates against the manual baseline.",
      Inches(0.5), Inches(4.7), W - Inches(1), Inches(0.72),
      13, MID, align=PP_ALIGN.CENTER)

    # CTA
    cta_w = Inches(3.4)
    cta_x = int((W - cta_w) / 2)
    R(s, cta_x, Inches(5.65), cta_w, Inches(0.6), BLUE)
    T(s, "Schedule a Pilot Run  →", cta_x + Inches(0.2), Inches(5.78),
      cta_w - Inches(0.4), Inches(0.36), 12.5, WHITE, True, PP_ALIGN.CENTER)

    # Three engineer characters at bottom
    chars = [
        (Inches(3.5),  BLUE,  "Instrumentation"),
        (Inches(6.67), TEAL,  "Piping"),
        (Inches(9.83), BLUE_D, "Management"),
    ]
    for cx, col, lbl in chars:
        person(s, cx, Inches(6.25), Inches(0.7), col, lbl)

    # Contact bar
    R(s, 0, H - Inches(0.72), W, Inches(0.72), SOFT)
    HL(s, 0, H - Inches(0.72), W, RULE, Pt(0.6))
    T(s, "prashanth.thipparthi@outlook.com",
      ML, H - Inches(0.55), Inches(7), Inches(0.36), 14, INK, True)
    T(s, "www.xyra-ai.com",
      W - Inches(3.5), H - Inches(0.55), Inches(3.28), Inches(0.36),
      12, BLUE, True, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════
#  BUILD
# ═══════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    for fn in [s01, s02, s03, s04, s05, s06, s07, s08, s09, s10, s11, s12, s13]:
        fn(prs)
    prs.save(str(OUT))
    print(f"✓  {OUT.name}")
    print(f"   {TOTAL} slides · White bg · Characters · Workflows · Real screenshots")

if __name__ == "__main__":
    build()
